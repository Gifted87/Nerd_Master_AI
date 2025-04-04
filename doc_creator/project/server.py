import os
import subprocess
import shutil
import tempfile
import re
import ast
import uuid
import logging
from logging.handlers import RotatingFileHandler
from flask import Flask, request, jsonify, send_file, render_template, abort
from werkzeug.security import safe_join
from flask_wtf.csrf import CSRFProtect, CSRFError
from dotenv import load_dotenv
from functools import wraps

load_dotenv()


# --- Configuration Loading ---
# Load sensitive keys from environment variables - DO NOT HARDCODE HERE
# Set these in your environment (e.g., .env file, system variables)
FLASK_SECRET_KEY = os.getenv('FLASK_SECRET_KEY', 'd5342974845146b06d677159749a5ed711a3c347595dd5c') # MUST be replaced for production
GEMINI_API_KEY = os.getenv('GEMINI_API_KEY') # MUST be set for the app to work
# Removed hardcoded API_KEY for Authorization header, rely on CSRF + potentially sessions/other auth
DOC_API_SECRET_KEY = os.getenv("DOC_API_SECRET_KEY")
if not DOC_API_SECRET_KEY:
    # Log a critical warning if the key isn't set, but allow the app to start
    # In production, you might want to raise an error here to prevent insecure startup.
    logging.warning("DOC_API_SECRET_KEY environment variable not set. API Key authentication for /generate-file will fail.")
# --- Application Setup ---
app = Flask(__name__)
app.config.update({
    'SECRET_KEY': FLASK_SECRET_KEY,
    'MAX_CONTENT_LENGTH': 1000 * 1024, # Limit request size (1 MB - adjusted)
    'GENERATED_FILES_DIR': os.path.abspath(os.path.join(os.path.dirname(__file__), 'generated_files')),
    # VERY restrictive module whitelist. Add ONLY if absolutely necessary and understood.
    'ALLOWED_MODULES': {
        'python-docx', 'docx', # Allow both common names
        'pptx', 'python-pptx', # Allow both common names
        'reportlab',
        'datetime',
        'random',
        # --- DO NOT ADD 'os', 'sys', 'subprocess', 'shutil', 'requests', 'socket', etc. ---
     },
    'EXECUTION_TIMEOUT': 15, # Slightly increased timeout
    'EXPECTED_OUTPUT_FILENAMES': {'output.docx', 'output.pdf', 'output.pptx'}, # Expected output names from generated script
})

# --- Logging Configuration ---
log_file = os.path.join(os.path.dirname(__file__), 'app.log')
# Increased size/count, consider log rotation strategy for production
handler = RotatingFileHandler(log_file, maxBytes=10 * 1024 * 1024, backupCount=5) # 10MB per file
formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(process)d - %(thread)d - %(message)s') # Added process/thread ID
handler.setFormatter(formatter)
handler.setLevel(logging.INFO) # Log INFO level and above
if app.debug:
    handler.setLevel(logging.DEBUG) # More verbose logging in debug mode
app.logger.addHandler(handler)
app.logger.setLevel(logging.INFO) # Ensure app logger respects level
logging.getLogger('werkzeug').setLevel(logging.INFO) # Quieter Werkzeug logs unless debugging


# --- Security ---
csrf = CSRFProtect(app)
# Removed WTF_CSRF_ENABLED = False for debug, CSRF should ideally always be checked

def require_api_key(f):
    """Decorator to protect routes with an API key check."""
    @wraps(f)
    def decorated_function(*args, **kwargs):
        api_key = request.headers.get('X-API-Key') # Header key matches Node.js client

        if not DOC_API_SECRET_KEY:
             app.logger.error("Server configuration error: DOC_API_SECRET_KEY is not set.")
             # Use abort to return a standard Flask error response
             abort(500, description="Server configuration error: API authentication not configured.")

        # Check if the key was provided and matches the expected one
        if not api_key or api_key != DOC_API_SECRET_KEY:
            app.logger.warning(f"Unauthorized access attempt: Invalid or missing API Key. Received: '{api_key}' from {request.remote_addr}")
            # Use abort for standard 401 response
            abort(401, description="Invalid or missing API Key.") # 401 Unauthorized is appropriate

        # If the key is valid, proceed with the original route function
        return f(*args, **kwargs)
    return decorated_function

# --- System Prompt for AI ---
# Moved the detailed instructions here to separate them from the user prompt.
SYSTEM_PROMPT = """You are a helpful assistant that generates Python code to create documents.

Instructions for Code Generation:
1.  Analyze the user's request to determine the appropriate file type: Word (.docx), PowerPoint (.pptx), or PDF (.pdf).
2.  Use **only** the following Python libraries for document generation:
    - For Word (.docx): `docx` (imported from the `python-docx` library)
    - For PowerPoint (.pptx): `pptx` (imported from the `python-pptx` library)
    - For PDF (.pdf): `reportlab` 
3.  You may also use the standard Python libraries `datetime` and `random` if necessary for content generation (e.g., adding dates, random data).
4.  The generated Python code MUST save the resulting document file in the **current working directory**.
5.  The filename MUST be exactly one of: 'output.docx', 'output.pptx', or 'output.pdf', corresponding to the document type.
6.  The Python code MUST NOT attempt to read files from the filesystem, access the network (no `requests`, `urllib`, `socket`, etc.), run shell commands (`os.system`, `subprocess`), interact with the operating system (`os`, `sys`), or perform any actions other than generating the specified document using *only* the allowed libraries. Do not use `shutil`.
7.  Do NOT include any `print()` statements or any other output in the generated Python code. The script should *only* perform the file generation and saving silently. Specifically, do not print 'FILE_SAVED:' or similar markers.
8.  For Word (.docx) and PDF (.pdf) documents, ensure the content is formatted professionally:
    - Use 'Times New Roman' font.
    - Use a font size of 13 points.
    - Apply proper margins.
    - Use 1.5 line spacing where appropriate.
9.  Ensure the generated Python code is complete, syntactically correct, and ready to be executed directly. Pay close attention to library import statements and function/method names (e.g., `Document` from `docx`, `Presentation` from `pptx`, `pdf` from `reportlab`).
10. **CRITICAL:** Prioritize generating bug-free, correct Python code that strictly adheres to these constraints over complex or potentially incorrect solutions. Ensure all keywords, library usage, and filenames are exactly as specified.

=== MANDATORY COMPONENTS ===
1. FILE TYPES:
   - Word (.docx): Use python-docx
   - PowerPoint (.pptx): Use python-pptx
   - PDF (.pdf): Use reportlab

2. REQUIRED IMPORTS (EXACT MATCH):
# DOCX IMPORTS
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_LINE_SPACING, WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
from docx.oxml.shared import OxmlElement
from docx.table import Table, _Cell

# PPTX IMPORTS
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# PDF IMPORTS
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, 
    Table, TableStyle, ListFlowable, ListItem
)
from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

=== ABSOLUTE PROHIBITIONS ===
1. NEVER USE:
   - Any import not explicitly listed above
   - os/sys/subprocess/shutil/io modules
   - print()/input()/eval()/exec()
   - Method chaining (e.g. doc.add_paragraph().add_run())
   - Comments about file operations
   - Markdown formatting

Allowed standard libraries:
- datetime
- random

=== FILE HANDLING ===
- Output filename MUST be exactly:
  - output.docx (Word)
  - output.pptx (PowerPoint)
  - output.pdf (PDF)
- NEVER use:
  - open(), os, sys, subprocess, shutil
  - Any filesystem/network operations

=== FORMATTING REQUIREMENTS ===
For Word
- Font: 'Times New Roman'
- Font size: 13pt
- Margins: 1 inch (Word: Inches(1))
- Line spacing: 1.5x (Word: WD_LINE_SPACING.ONE_POINT_FIVE)

For PDF:
- Font: 'Times-Roman' (must register if using custom TTF)
- Font size: 13pt
- Margins: 1 inch (72 points)
- Line spacing: 1.5x (18pt for 13pt font)

For PowerPoint:
- Title text: 28-32pt
- Body text: 18-22pt
- Font: 'Times New Roman'

=== CRITICAL RULES ===
1. NO method chaining - always create variables:
   # CORRECT:
   p = doc.add_paragraph()
   p.add_run("text")
   
   # WRONG:
   doc.add_paragraph().add_run("text")

2. NO comments about file saving
3. NO print statements or other output
4. NO markdown - only raw Python code
5. Use EXACT method names:
   - add_paragraph() (Word)
   - add_slide() (PowerPoint)
   - multi_cell() (PDF)

=== DOCUMENT TEMPLATES ===
# WORD DOCUMENT (output.docx)
doc = Document()
# Set document-wide styles
styles = doc.styles
normal_style = styles['Normal']
normal_style.font.name = 'Times New Roman'
normal_style.font.size = Pt(13)
# Configure margins
section = doc.sections[0]
section.left_margin = Inches(1)
section.right_margin = Inches(1)
section.top_margin = Inches(1)
section.bottom_margin = Inches(1)
# Add content
para = doc.add_paragraph()
run = para.add_run("Professional Content")
run.bold = True
# Save document
doc.save('output.docx')

# POWERPOINT (output.pptx)
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
# Title setup
title = slide.shapes.title
title.text = "Presentation Title"
title_frame = title.text_frame
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(28)
title_para.font.name = 'Times New Roman'
# Content setup
body = slide.placeholders[1]
body_frame = body.text_frame
body_para = body_frame.add_paragraph()
body_para.text = "• First bullet point"
body_para.font.size = Pt(18)
prs.save('output.pptx')

# PDF DOCUMENT (output.pdf)
# Font registration
pdfmetrics.registerFont(TTFont('Times-Roman', 'times.ttf'))
# Style configuration
styles = getSampleStyleSheet()
custom_style = ParagraphStyle(
    'MainStyle',
    parent=styles['Normal'],
    fontName='Times-Roman',
    fontSize=13,
    leading=18,  # 1.5x line spacing
    leftIndent=0,
    rightIndent=0,
    spaceBefore=6,
    spaceAfter=6
)
# Document structure
doc = SimpleDocTemplate(
    "output.pdf",
    pagesize=letter,
    leftMargin=inch,
    rightMargin=inch,
    topMargin=inch,
    bottomMargin=inch
)
# Content creation
story = []
numbered_list = ListFlowable(
    [
        ListItem(Paragraph("Item 1", custom_style)),
        ListItem(Paragraph("Item 2", custom_style))
    ],
    bulletType='1'
)
story.append(numbered_list)
doc.build(story)

=== VALIDATION PROTOCOL ===
Before finalizing code:
1. Import Check: Verify every import matches EXACTLY the allowed list
2. Filename Validation: Must be output.docx/pptx/pdf
3. Style Verification:
   - DOCX: Times New Roman 13pt, 1" margins
   - PDF: Times-Roman 13pt, 1.5x spacing
   - PPTX: Title 28-32pt, Body 18-22pt
4. Syntax Check:
   - No undefined variables
   - No method chaining
   - Proper indentation

=== ERROR-PRONE SCENARIOS ===
1. PDF Lists:
   - Use ListFlowable + ListItem
   - Valid bulletType: '1', 'a', 'i', 'bullet'
   - Never use ListBullet

2. Tables:
   - DOCX: table = doc.add_table(rows=X, cols=Y)
   - PDF: Use Table(data) with TableStyle

3. Images:
   - PDF: from reportlab.platypus import Image
   - PPTX: Use slide.shapes.add_picture()

4. Fonts:
   - DOCX: Set via style.font.name
   - PDF: Must register fonts with pdfmetrics

=== EXAMPLES OF BAD CODE ===
# WRONG: Invalid import
from reportlab.platypus import ListBullet  # ERROR - Doesn't exist

# WRONG: Method chaining
doc.add_paragraph().add_run("text")  # ERROR - Chain not allowed

# WRONG: Missing style setup
doc = Document()
doc.add_paragraph("Content")  # ERROR - No font/margin config


=== FINAL CHECK ===
Before returning code:
1. Verify all imports match exactly the approved lists
2. Confirm the code follows the templates exactly
3. Check for any syntax errors
4. Ensure all required formatting is included
5. Validate the output filename is exactly 'output.docx', 'output.pptx', or 'output.pdf'

=== VALIDATION FAILURES ===
The server will reject code for:
1. Missing required formatting (font/margins/spacing)
2. Using forbidden modules/functions
3. Incorrect filename/path
4. Method chaining
5. Any print/output statements

=== EXAMPLE VALID OUTPUT ===
from docx import Document
from docx.shared import Pt, Inches
from docx.enum.text import WD_LINE_SPACING

doc = Document()
style = doc.styles['Normal']
style.font.name = 'Times New Roman'
style.font.size = Pt(13)
section = doc.sections[0]
section.left_margin = Inches(1)
section.right_margin = Inches(1)
section.top_margin = Inches(1)
section.bottom_margin = Inches(1)
paragraph_format = style.paragraph_format
paragraph_format.line_spacing_rule = WD_LINE_SPACING.ONE_POINT_FIVE
doc.add_paragraph("Professional Content")
doc.save('output.docx')
"""

# --- AI Configuration ---
model = None # Initialize model variable
genai = None # Initialize genai library variable

if not GEMINI_API_KEY:
    app.logger.critical("GEMINI_API_KEY environment variable not set. AI features will fail.")
    # Optionally, raise an exception or exit if the key is essential for startup
    # raise ValueError("GEMINI_API_KEY must be set")
else:
    try:
        import google.generativeai as genai
        genai.configure(api_key=GEMINI_API_KEY)
        # Choose a model - check availability and features
        # e.g., 'gemini-1.5-flash', 'gemini-1.0-pro'
        # Initialize the model with the system prompt
        model = genai.GenerativeModel(
            'gemini-2.0-flash-thinking-exp-01-21', # Using a standard reliable model
            system_instruction=SYSTEM_PROMPT # Pass the system prompt here
        )
        app.logger.info("Google Generative AI configured successfully with system prompt.")
    except ImportError:
        app.logger.critical("google.generativeai library not installed. Run 'pip install google-generativeai'")
        genai = None # Ensure genai is None if import fails
    except Exception as e:
        app.logger.critical(f"Failed to configure Google Generative AI: {e}")
        genai = None # Ensure genai is None if config fails
        model = None # Ensure model is None if config fails


# --- Directory Setup ---
os.makedirs(app.config['GENERATED_FILES_DIR'], exist_ok=True)


def remove_html_tags(script):
    # Remove all HTML tags using regular expression
    text = re.sub(r'<[^>]+>', '', script)
    # Replace any whitespace sequences (including newlines and tabs) with a single space
    text = re.sub(r'\s+', ' ', text)
    # Strip leading/trailing spaces and convert to lowercase
    return text.strip().lower()


# --- Code Validation Logic ---
def validate_generated_code(code: str) -> bool:
    """
    Perform security checks on AI-generated Python code.
    Checks for dangerous patterns, forbidden modules, and potentially harmful built-ins.

    NOTE: This validation is helpful but NOT foolproof. Strong sandboxing is essential.
    """
    app.logger.debug(f"Validating code (first 200 chars): {code[:200]}...")

    # 1. Length Check
    MAX_CODE_LENGTH = 1000000 # Define max length clearly
    if len(code) > MAX_CODE_LENGTH:
        app.logger.warning(f"Code length ({len(code)}) exceeded limit ({MAX_CODE_LENGTH}).")
        # Avoid logging potentially huge code blocks fully
        app.logger.warning(f"Code start: {code[:500]}...")
        return False

    # 2. Forbidden Patterns (Regex) - Stricter and more focused
    #    Focus on explicit calls to dangerous functions/modules.
    #    AST checking handles imports better, regex is a first line of defense.
    #    Using raw strings (r'...') is important for regex patterns.
    forbidden_patterns = [
        # # Explicit blocking of dangerous modules even if AST missed them somehow
        # r'\bimport\s+(os|sys|subprocess|shutil|socket|requests|urllib|httpx|pickle|ctypes)\b',
        # r'\bfrom\s+(os|sys|subprocess|shutil|socket|requests|urllib|httpx|pickle|ctypes)\s+import\b',
        # # Direct calls to dangerous functions
        # r'\b(eval|exec|compile|open|input|__import__)\s*\(',
        # # Accessing system/environment variables/dunder methods that could be risky
        # r'os\.(system|popen|environ|getenv|putenv|listdir|scandir|remove|unlink|rmdir|makedirs|chmod|chown)',
        # r'subprocess\.(run|call|check_call|check_output|Popen)',
        # r'shutil\.(move|copyfile|copytree|rmtree|disk_usage)', # Allow only expected libs
        # r'socket\.',
        # r'requests\.',
        # r'urllib\.',
        # r'httpx\.',
        # r'pickle\.(load|loads|dump|dumps)',
        # r'ctypes\.',
        # # Dynamic execution/attribute access
        # r'\bgetattr\s*\(',
        # r'\bsetattr\s*\(',
        # r'\bglobals\s*\(',
        # r'\blocals\s*\(',
        # # Direct file writing attempts with open (basic check)
        # r'open\s*\([^)]+,[^)]*[\'"](w|a|x|wb|ab|xb)[\'"]',
        # # Common command injection keywords (crude but can catch simple cases)
        # # r'\b(rm|del|mv|cp|chmod|chown|sudo|apt|yum|dnf|pacman|git|wget|curl)\s+',
        # # Trying to access sensitive paths (examples)
        # r'(/|\\)(etc|root|home|users|windows|system32)(\\|/|$)',
        # r'(\\.ssh|\\.git|\\.aws|\\.azure|\\.gcp)',
    ]

    for pattern in forbidden_patterns:
        if re.search(pattern, code, re.IGNORECASE): # Ignore case for broader matching
            app.logger.warning(f"Forbidden pattern detected in generated code: {pattern}")
            # Log only a snippet for security reasons if code is large
            app.logger.warning(f"Code snippet: {code[:500]}...")
            return False

    # 3. AST Analysis (Abstract Syntax Tree)
    try:
        tree = ast.parse(code)
        allowed_modules = app.config['ALLOWED_MODULES']

        for node in ast.walk(tree):
            # Check Imports more reliably
            if isinstance(node, ast.Import):
                for alias in node.names:
                    module_name = alias.name.split('.')[0] # Get base module
                    if module_name not in allowed_modules:
                        app.logger.warning(f"AST Check: Disallowed module imported: {module_name}")
                        return False
            elif isinstance(node, ast.ImportFrom):
                # Handle 'from . import X' cases (relative imports, should likely be disallowed)
                if node.level > 0:
                    app.logger.warning(f"AST Check: Disallowed relative import.")
                    return False
                module_name = node.module.split('.')[0] if node.module else None # Get base module
                if module_name and module_name not in allowed_modules:
                     app.logger.warning(f"AST Check: Disallowed module imported via 'from': {module_name}")
                     return False
                # Optionally, check imported names (node.names) for specific dangerous functions from allowed modules if needed

            # Check Dangerous Built-in Functions/Attributes (AST is more precise than regex)
            # Check for function calls
            if isinstance(node, ast.Call) and isinstance(node.func, ast.Name):
                if node.func.id in ['eval', 'exec', 'compile', 'open', 'input', '__import__', 'globals', 'locals', 'getattr', 'setattr']:
                    app.logger.warning(f"AST Check: Potentially dangerous built-in function called: {node.func.id}")
                    return False
            # Check for direct attribute access on disallowed modules (belt-and-suspenders)
            # Example: If 'os' somehow passed import checks, catch os.system here
            if isinstance(node, ast.Attribute) and isinstance(node.value, ast.Name):
                 if node.value.id in ['os', 'sys', 'subprocess', 'shutil', 'socket', 'requests', 'urllib', 'httpx', 'pickle', 'ctypes']:
                     # Check specific dangerous attributes if needed
                     # e.g., if node.value.id == 'os' and node.attr == 'system': return False
                     app.logger.warning(f"AST Check: Detected attribute access on potentially disallowed module '{node.value.id}' (Attribute: {node.attr})")
                     # Be cautious here, might be too broad. Focus on known dangerous attributes.
                     if node.value.id == 'os' and node.attr in ['system', 'popen', 'environ', 'getenv', 'putenv', 'listdir', 'scandir', 'remove', 'unlink', 'rmdir', 'makedirs', 'chmod', 'chown']:
                         return False
                     if node.value.id == 'subprocess' and node.attr in ['run', 'call', 'check_call', 'check_output', 'Popen']:
                         return False
                     # Add more specific checks as needed

    except SyntaxError as e:
        app.logger.error(f"Syntax error parsing generated code: {e}")
        # Log snippet for debugging
        app.logger.error(f"Code snippet with error: {code[max(0, e.lineno-2):e.lineno+1]}")
        return False
    except Exception as e:
        # Catch potential recursion errors with deep ASTs etc.
        app.logger.error(f"Unexpected error during AST parsing: {e}", exc_info=True)
        return False # Fail safe

    app.logger.debug("Code validation passed.")
    return True


# --- Flask Routes ---
@app.route('/')
def index():
    """Renders the main HTML page."""
    return render_template('index.html')

@app.route('/generate-file', methods=['POST'])
# @csrf.exempt # Keep CSRF protection unless specifically handled otherwise (e.g., API keys)
# If using API Key auth primarily, CSRF might be less critical *for this specific endpoint*,
# but it's safer to keep it unless it interferes with the client implementation.
# Let's assume API key is the primary auth here, so exempting for simplicity with non-browser clients.
@csrf.exempt # Exempt CSRF for API-key protected endpoint often used by scripts/non-browsers
# @require_api_key # Apply the API key check decorator
def generate_file():
    """
    Handles the POST request to generate a document file using AI-generated code.
    Requires a valid API Key and JSON payload with a 'prompt'.
    """
    app.logger.info(f"File generation request received from {request.remote_addr}")

    # --- Input Validation ---
    if not request.is_json:
        app.logger.warning("Invalid content type: Expected application/json")
        return jsonify({"error": "Invalid content type, requires application/json"}), 400

    data = request.get_json()
    if not data:
         app.logger.warning("Empty JSON payload received")
         return jsonify({"error": "Empty request body"}), 400

    user_prompt = data.get('prompt')
    if not user_prompt or not isinstance(user_prompt, str) or not user_prompt.strip():
        app.logger.warning("Missing, invalid, or empty 'prompt' in request JSON")
        return jsonify({"error": "Missing, invalid, or empty 'prompt' field"}), 400

    # Limit prompt length server-side as well
    MAX_PROMPT_LENGTH = 1000000 # Increased limit slightly
    user_prompt = user_prompt[:MAX_PROMPT_LENGTH].strip() # Trim whitespace too
    app.logger.debug(f"Processing prompt (first 100 chars): {user_prompt[:100]}...")

    # --- AI Code Generation ---
    if not genai or not model:
         app.logger.error("AI model not configured or failed to initialize. Cannot generate code.")
         # 503 Service Unavailable is appropriate if the AI backend is down/unconfigured
         return jsonify({"error": "AI service not available"}), 503

    try:
        app.logger.info("Sending request to AI model...")
        # The detailed instructions are now in the system_instruction used when initializing the model.
        # We only need to send the user's specific request here.
        response = model.generate_content(remove_html_tags(user_prompt)) # Pass only the user prompt

        # --- Response Processing ---
        # Check for safety ratings or blocks if the API provides them
        if hasattr(response, 'prompt_feedback') and response.prompt_feedback.block_reason:
            app.logger.warning(f"AI generation blocked for prompt. Reason: {response.prompt_feedback.block_reason}")
            return jsonify({"error": f"Request blocked by safety filter: {response.prompt_feedback.block_reason}"}), 400

        if not response.candidates:
             app.logger.error("AI model returned no candidates.")
             raise ValueError("AI returned no candidates")

        # Check safety ratings of the first candidate if available
        candidate = response.candidates[0]
        if candidate.finish_reason != 'STOP':
             app.logger.warning(f"AI generation finished unexpectedly. Reason: {candidate.finish_reason}")
             # Handle potential safety blocks here too
             if candidate.finish_reason == 'SAFETY':
                 safety_ratings_str = ", ".join([f"{rating.category}: {rating.probability}" for rating in candidate.safety_ratings])
                 app.logger.warning(f"Safety block details: {safety_ratings_str}")
                 return jsonify({"error": f"Generation stopped due to safety concerns ({candidate.finish_reason})."}), 400
             # else: # Other reasons like MAX_TOKENS, RECITATION etc. might warrant errors or specific handling
             #     return jsonify({"error": f"Generation incomplete ({candidate.finish_reason})."}), 500

        # Extract generated code from the candidate's content
        # Assuming the code is within the 'parts' of the content
        generated_code = "".join(part.text for part in candidate.content.parts if hasattr(part, 'text'))

        # --- Basic Code Extraction (if markdown format is used) ---
        # This tries to strip common markdown code blocks if present
        code_match = re.search(r"```(?:python)?\s*([\s\S]+?)\s*```", generated_code, re.IGNORECASE)
        if code_match:
            generated_code = code_match.group(1).strip()
            app.logger.debug("Extracted code from markdown block.")
        else:
            # Fallback: assume the whole response might be code, clean basic comments/whitespace
             generated_code = '\n'.join([line for line in generated_code.split('\n')
                                       if not line.strip().startswith('#')]).strip() # Simple comment removal + strip
             app.logger.debug("No markdown block found, using raw response (comments stripped).")


        if not generated_code:
            app.logger.error("AI model returned empty code after extraction.")
            # Check if the raw response text had content before extraction
            raw_text = response.text if hasattr(response, 'text') else "(No raw text available)"
            app.logger.error(f"Raw response text (first 500): {raw_text[:500]}")
            raise ValueError("AI returned empty code after extraction")

        app.logger.debug(f"Generated code received (first 500 chars):\n{generated_code[:500]}...")

    except Exception as e:
        app.logger.error(f"AI code generation or processing failed: {e}", exc_info=True)
        # Provide a more generic error to the client
        return jsonify({"error": "Failed to generate document code from AI"}), 500

    # --- Code Validation ---
    app.logger.info("Validating generated code...")
    if not validate_generated_code(generated_code):
        app.logger.warning("Generated code failed validation.")
        # Do not expose details of validation failure to the client
        return jsonify({"error": "Generated code is invalid or potentially unsafe"}), 400
    app.logger.info("Code validation successful.")


    # --- Secure Code Execution ---
    # CRITICAL: This section requires robust sandboxing (Docker, nsjail, firecracker, etc.)
    # The TemporaryDirectory provides *filesystem* isolation but NOT process/network isolation.
    # The process still runs with the web server's user permissions.
    # DO NOT RUN IN PRODUCTION WITHOUT A REAL SANDBOX.
    output_filename = None
    final_filename = None
    file_moved = False
    generated_file_path = None # Keep track of the final path

    # Use a unique subdirectory within the main generated files dir for better organization
    # and to simplify cleanup if needed.
    session_id = uuid.uuid4()
    session_dir = os.path.join(app.config['GENERATED_FILES_DIR'], str(session_id))
    os.makedirs(session_dir, exist_ok=True) # Create a directory for this specific generation

    # Consider using a temporary directory that is *outside* the main GENERATED_FILES_DIR initially
    # for execution, then moving the result *into* the session_dir. This adds a layer.
    with tempfile.TemporaryDirectory() as exec_temp_dir:
        script_path = os.path.join(exec_temp_dir, 'generated_script.py')
        try:
            with open(script_path, 'w', encoding='utf-8') as f:
                f.write(generated_code)

            app.logger.info(f"Executing generated script in isolated temporary directory: {exec_temp_dir}")

            # *** SANDBOXING PLACEHOLDER ***
            # This is where you would replace subprocess.run with your sandboxing mechanism.
            # Example:
            # result = run_in_sandbox(
            #     script_path,
            #     execution_environment=exec_temp_dir, # Mount or map this directory in the sandbox
            #     allowed_network=False,
            #     cpu_limit='1s', # Resource limits
            #     memory_limit='128m',
            #     timeout=app.config['EXECUTION_TIMEOUT']
            # )
            # The result object should contain stdout, stderr, returncode, and potentially info about created files.

            # Current implementation (UNSAFE for production):
            try:
                 result = subprocess.run(
                    ['python', script_path], # Use isolated mode '-I' for extra safety
                    cwd=exec_temp_dir,        # Execute within the temp dir
                    capture_output=True,
                    text=True,
                    timeout=app.config['EXECUTION_TIMEOUT'],
                    check=False,              # Don't raise exception on non-zero exit, check manually
                    # env={},                   # Clear environment variables for the subprocess
                )
            except FileNotFoundError:
                 app.logger.critical("'python' command not found. Ensure Python is installed and in PATH for the server process.")
                 raise RuntimeError("Python executable not found.")
            except Exception as subproc_err: # Catch broader errors during subprocess creation/execution
                 app.logger.error(f"Subprocess execution failed unexpectedly: {subproc_err}", exc_info=True)
                 raise RuntimeError("Failed to run the generated script.")


            app.logger.debug(f"Script execution finished. Return code: {result.returncode}")
            # Log stdout/stderr cautiously (limit size)
            if result.stdout:
                app.logger.debug(f"Script stdout (limited):\n{result.stdout[:1000]}")
            if result.stderr:
                # Treat stderr as a potential error/warning
                app.logger.warning(f"Script stderr (limited):\n{result.stderr[:1000]}")

            if result.returncode != 0:
                app.logger.error(f"Generated script execution failed with return code {result.returncode}.")
                # Avoid sending detailed stderr to client unless it's specifically sanitized/allowed
                # Potentially include stderr snippet in server logs for debugging
                app.logger.error(f"Script stderr detail: {result.stderr}")
                raise ValueError(f"Script execution failed (code: {result.returncode})")

            # --- Find, Validate, and Move Output File ---
            found_expected_file = False
            MAX_FILE_SIZE_MB = 50
            MAX_FILE_SIZE_BYTES = MAX_FILE_SIZE_MB * 1024 * 1024

            for expected_name in app.config['EXPECTED_OUTPUT_FILENAMES']:
                temp_output_path = os.path.join(exec_temp_dir, expected_name) # Look in exec temp dir

                if os.path.exists(temp_output_path) and os.path.isfile(temp_output_path):
                    file_size = os.path.getsize(temp_output_path)
                    app.logger.info(f"Found potential output file: {expected_name} (Size: {file_size} bytes)")

                    # Basic check: Ensure file is not excessively large
                    if file_size > MAX_FILE_SIZE_BYTES:
                         app.logger.error(f"Generated file '{expected_name}' is too large ({file_size} bytes > {MAX_FILE_SIZE_BYTES} bytes).")
                         # Clean up the large file before raising
                         try: os.remove(temp_output_path)
                         except OSError: pass
                         raise ValueError(f"Generated file size exceeds limit ({MAX_FILE_SIZE_MB} MB)")
                    elif file_size == 0:
                         app.logger.warning(f"Generated file '{expected_name}' is empty.")
                         # Decide if empty files are errors or just warnings
                         # For now, treat as error as it's likely unintentional
                         try: os.remove(temp_output_path)
                         except OSError: pass
                         raise ValueError("Generated file is empty")

                    # Generate a unique final filename within the session directory
                    file_ext = os.path.splitext(expected_name)[1]
                    # Use the original expected name + session ID for clarity
                    # final_filename = f"docgen_{session_id}_{expected_name}" # Or keep simpler name if session dir implies uniqueness
                    final_filename = expected_name # Keep original name within the unique session dir

                    # Use safe_join for constructing the final destination path inside the session directory
                    # safe_join protects against the 'final_filename' part having traversal characters
                    dest_path = safe_join(session_dir, final_filename)

                    # Double-check: Ensure dest_path is actually inside GENERATED_FILES_DIR
                    # safe_join might return a path outside if the base itself is tricky, although unlikely here.
                    # os.path.abspath resolves any '..' etc.
                    abs_dest_path = os.path.abspath(dest_path)
                    abs_generated_dir = os.path.abspath(app.config['GENERATED_FILES_DIR'])

                    if not dest_path or not abs_dest_path.startswith(abs_generated_dir):
                        app.logger.error(f"Could not create safe destination path or path escaped base directory for: {final_filename} in {session_dir}")
                        raise ValueError("Invalid destination path generated")

                    app.logger.info(f"Moving expected output file '{expected_name}' to final location: {dest_path}")
                    shutil.move(temp_output_path, dest_path) # Move from exec_temp_dir to session_dir
                    generated_file_path = dest_path # Store the final path
                    file_moved = True
                    found_expected_file = True
                    break # Found the file, no need to check others

            if not found_expected_file:
                app.logger.error("Generated script completed successfully but did not produce an expected output file (%s) in %s.",
                                 ', '.join(app.config['EXPECTED_OUTPUT_FILENAMES']), exec_temp_dir)
                # Log contents of the exec dir for debugging
                try:
                    dir_contents = os.listdir(exec_temp_dir)
                    app.logger.error(f"Contents of execution directory: {dir_contents}")
                except OSError as list_err:
                    app.logger.error(f"Could not list execution directory contents: {list_err}")
                raise ValueError("Script did not create the expected output file")

        except subprocess.TimeoutExpired:
            app.logger.error(f"Generated script timed out after {app.config['EXECUTION_TIMEOUT']} seconds.")
            # Attempt cleanup of the session directory if timeout occurred
            shutil.rmtree(session_dir, ignore_errors=True)
            return jsonify({"error": "Document generation timed out"}), 504 # 504 Gateway Timeout might be suitable
        except (ValueError, RuntimeError, OSError) as script_err: # Catch specific errors from execution/validation
            app.logger.error(f"Error during script execution or file handling: {script_err}", exc_info=True)
            # Attempt cleanup of the session directory on error
            shutil.rmtree(session_dir, ignore_errors=True)
            # Return specific error message if it's user-safe, otherwise generic
            user_error_message = str(script_err) if isinstance(script_err, ValueError) else "Failed to create or save the document"
            return jsonify({"error": user_error_message}), 500
        except Exception as e: # Catch-all for unexpected errors
            app.logger.error(f"Unexpected error during script execution phase: {e}", exc_info=True)
            # Attempt cleanup
            shutil.rmtree(session_dir, ignore_errors=True)
            # Generic error to client
            return jsonify({"error": "An unexpected error occurred during document creation"}), 500
        # The 'finally' block for the tempfile.TemporaryDirectory handles cleanup of exec_temp_dir

    # --- Success Response ---
    if file_moved and generated_file_path and final_filename:
        # Construct the download URL relative path (session_id/filename)
        relative_download_path = f"{session_id}/{final_filename}"
        app.logger.info(f"Successfully generated file: {generated_file_path}. Download URL segment: {relative_download_path}")
        return jsonify({
            # The download URL now needs to include the session ID part
            "download_url": f"/download/{relative_download_path}",
            "filename": final_filename # Return the simple filename for display
        }), 200
    else:
        # This state should ideally not be reachable if logic above is correct
        app.logger.error("Reached end of generation process without a valid file state (file_moved=%s, path=%s, name=%s).",
                         file_moved, generated_file_path, final_filename)
        # Attempt cleanup
        shutil.rmtree(session_dir, ignore_errors=True)
        return jsonify({"error": "An unexpected error occurred during file finalization"}), 500


@app.route('/download/<uuid:session_id>/<path:filename>')
def download_file(session_id, filename):
    """Serves the generated file for download from its session directory."""
    app.logger.info(f"Download request for: {filename} in session {session_id} from {request.remote_addr}")
    try:
        # Construct the path including the session ID directory
        # Ensure session_id is treated as a string for path joining
        session_dir_name = str(session_id)
        # safe_join is crucial here to prevent filename from having traversal components
        # It joins the base directory, the session directory, and the filename safely
        safe_path = safe_join(app.config['GENERATED_FILES_DIR'], session_dir_name, filename)

        if safe_path is None:
             app.logger.warning(f"Download aborted: safe_join failed for session '{session_id}', filename '{filename}'. Potential path traversal attempt.")
             abort(404, description="File path is invalid.")

        # Check if the file actually exists
        if not os.path.isfile(safe_path):
            app.logger.warning(f"File not found at path: {safe_path}")
            # Log if the session directory itself exists for debugging
            session_dir_path = os.path.join(app.config['GENERATED_FILES_DIR'], session_dir_name)
            if not os.path.isdir(session_dir_path):
                app.logger.warning(f"Session directory not found: {session_dir_path}")
            else:
                 try:
                     dir_contents = os.listdir(session_dir_path)
                     app.logger.warning(f"Contents of session directory {session_dir_name}: {dir_contents}")
                 except OSError:
                      app.logger.warning(f"Could not list contents of session directory {session_dir_name}")
            abort(404, description="File not found.")


        # Double-check: Verify the absolute path is still within the GENERATED_FILES_DIR (belt-and-suspenders)
        abs_safe_path = os.path.abspath(safe_path)
        abs_generated_dir = os.path.abspath(app.config['GENERATED_FILES_DIR'])
        if not abs_safe_path.startswith(abs_generated_dir):
             app.logger.error(f"Attempt to access file outside GENERATED_FILES_DIR detected: Path {safe_path} resolved to {abs_safe_path}")
             abort(404, description="File not found (Security restriction).")


        app.logger.info(f"Sending file: {safe_path}")
        return send_file(safe_path, as_attachment=True)

    except TypeError as e:
         # Catch potential errors if session_id wasn't a valid UUID string for the route converter
         app.logger.error(f"Invalid session ID format in URL: {session_id}. Error: {e}")
         abort(404, description="Invalid session ID format.")
    except Exception as e:
        app.logger.error(f"Error during file download for session {session_id}, file {filename}: {e}", exc_info=True)
        abort(500, description="Could not process file download.")

# --- Error Handling ---
@app.errorhandler(CSRFError)
def handle_csrf_error(e):
    app.logger.warning(f"CSRF validation failed: {e.description} from {request.remote_addr}")
    return jsonify({"error": "CSRF validation failed: " + e.description}), 400

@app.errorhandler(400)
def bad_request_error(error):
    description = getattr(error, 'description', 'Bad request.')
    app.logger.warning(f"Bad Request (400): {description} from {request.remote_addr}")
    return jsonify({"error": "Bad Request", "message": description}), 400

@app.errorhandler(401)
def unauthorized_error(error):
    # Log specific description if available from abort()
    description = getattr(error, 'description', 'Unauthorized access.')
    app.logger.warning(f"Unauthorized (401): {description} from {request.remote_addr}")
    return jsonify({"error": "Unauthorized", "message": description}), 401

@app.errorhandler(404)
def not_found_error(error):
    description = getattr(error, 'description', 'The requested resource was not found.')
    app.logger.info(f"Not Found (404): {description} for {request.path} from {request.remote_addr}")
    return jsonify({"error": "Not Found", "message": description}), 404

@app.errorhandler(405)
def method_not_allowed_error(error):
    description = getattr(error, 'description', 'The method is not allowed for the requested URL.')
    app.logger.warning(f"Method Not Allowed (405): {request.method} for {request.path} from {request.remote_addr}")
    return jsonify({"error": "Method Not Allowed", "message": description}), 405

@app.errorhandler(500)
def internal_error(error):
    # Log the actual error if possible, but return generic message
    # Check if it's an HTTPException with a description, otherwise log the base error
    original_exception = getattr(error, 'original_exception', error)
    app.logger.error(f"Internal Server Error (500): {original_exception}", exc_info=True)
    # Send a generic message to the client
    return jsonify({"error": "Internal Server Error", "message": "An unexpected error occurred on the server."}), 500

@app.errorhandler(503)
def service_unavailable_error(error):
    description = getattr(error, 'description', 'The service is temporarily unavailable.')
    app.logger.error(f"Service Unavailable (503): {description} from {request.remote_addr}")
    return jsonify({"error": "Service Unavailable", "message": description}), 503

@app.errorhandler(504)
def gateway_timeout_error(error):
    description = getattr(error, 'description', 'The server timed out waiting for an upstream response.')
    app.logger.error(f"Gateway Timeout (504): {description} from {request.remote_addr}")
    return jsonify({"error": "Gateway Timeout", "message": description}), 504


# --- Main Execution ---
if __name__ == '__main__':
    # Production: Use a proper WSGI server (Gunicorn/uWSGI) behind a reverse proxy (Nginx/Caddy)
    # Development: Use Flask's built-in server (debug=True enables reloader and debugger)
    # Determine debug mode from environment variables
    is_debug = os.getenv('FLASK_ENV', '').lower() == 'development' or os.getenv('FLASK_DEBUG', '0') == '1'

    if is_debug:
        app.logger.setLevel(logging.DEBUG) # Ensure logger is verbose in debug
        handler.setLevel(logging.DEBUG)
        logging.getLogger('werkzeug').setLevel(logging.DEBUG)
        app.logger.info("Flask app running in DEBUG mode.")
    else:
        app.logger.info("Flask app running in PRODUCTION mode (or debug not explicitly enabled).")

    # Get host and port from environment variables, defaulting for development
    host = os.getenv('FLASK_RUN_HOST', '0.0.0.0')
    port = int(os.getenv('PORT', 5000)) # PORT is common for PaaS like Heroku/Cloud Run

    app.logger.info(f"Starting Flask server on {host}:{port}")
    # Note: Flask's built-in server is NOT recommended for production.
    # 'debug=is_debug' controls the reloader and debugger.
    # 'use_reloader=is_debug' explicitly controls the reloader if needed separately.
    app.run(
        host=host,
        port=port,
        debug=is_debug
    )
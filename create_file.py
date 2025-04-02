```python
import os
from docx import Document
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

def generate_file(user_prompt):
    """Generates a file based on the user prompt."""

    if "Word" in user_prompt or ".docx" in user_prompt.lower():
        return create_word_document(user_prompt)
    elif "PPTX" in user_prompt or ".pptx" in user_prompt.lower():
        return create_powerpoint_presentation(user_prompt)
    elif "PDF" in user_prompt or ".pdf" in user_prompt.lower():
        return create_pdf_document(user_prompt)
    else:
        return None

def create_word_document(user_prompt):
    """Creates a Word document based on the user prompt."""
    document = Document()
    document.add_heading('Meeting Minutes', level=1)
    document.add_paragraph('Meeting Date: 2024-01-26')
    document.add_paragraph('Attendees:')
    document.add_paragraph('- John Doe')
    document.add_paragraph('- Jane Smith')
    document.add_paragraph('Agenda Items:')
    document.add_paragraph('1. Project Status Update')
    document.add_paragraph('2. Budget Review')
    document.add_paragraph('3. Next Steps')
    filename = "meeting_minutes.docx"
    document.save(filename)
    return filename

def create_powerpoint_presentation(user_prompt):
    """Creates a PowerPoint presentation based on the user prompt."""
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.shapes.placeholders[1]
    title.text = 'Project Proposal'
    subtitle.text = 'Prepared by John Doe'

    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    title.text = 'Key Objectives'
    tf = body.text_frame
    tf.text = 'Increase market share by 15%'
    p = tf.add_paragraph()
    p.text = 'Improve customer satisfaction'
    p = tf.add_paragraph()
    p.text = 'Reduce operational costs'
    filename = "project_proposal.pptx"
    prs.save(filename)
    return filename

def create_pdf_document(user_prompt):
    """Creates a PDF document based on the user prompt."""
    filename = "report.pdf"
    c = canvas.Canvas(filename, pagesize=letter)
    c.drawString(100, 750, "Monthly Sales Report")
    c.drawString(100, 730, "Date: 2024-01-26")
    c.drawString(100, 700, "Sales Figures:")
    c.drawString(100, 680, "Product A: $10,000")
    c.drawString(100, 660, "Product B: $15,000")
    c.drawString(100, 640, "Product C: $8,000")
    c.save()
    return filename

user_prompt = input("Enter the desired file type and content request: ")
filename = generate_file(user_prompt)

if filename:
    print(f"FILE_SAVED: {filename}")
```